import os
import sys
import tempfile
import shutil
import threading
import tkinter as tk
from tkinter import filedialog, messagebox, ttk

# 第三方库导入
try:
    import win32com.client
    from pptx import Presentation as PPTXPresentation
    from pptx.util import Emu
    from pptx.dml.color import RGBColor
    from PIL import Image
except ImportError as e:
    missing = str(e).split()[-1].strip("'")
    messagebox.showerror("缺少依赖库",
                         f"请先安装所需库：\n\npip install {missing}\n\n"
                         f"如果缺少多个库，请依次安装：pywin32, python-pptx, Pillow")
    sys.exit(1)

# ========== 核心转换函数（支持DPI选择与压缩） ==========
def ppt_to_high_res_ppt(input_ppt_path, output_pptx_path,
                        dpi_target, enable_compress, compress_quality,
                        progress_callback, status_callback):
    """
    将PPT/PPTX文件转换为高分辨率图片型PPT
    :param input_ppt_path: 输入PPT/PPTX文件路径
    :param output_pptx_path: 输出PPTX文件路径
    :param dpi_target: 导出图片的目标DPI（整数）
    :param enable_compress: 是否启用图片压缩
    :param compress_quality: JPEG压缩质量（1-100）
    :param progress_callback: 进度回调
    :param status_callback: 状态回调
    :return: (success, error_message)
    """
    temp_dir = None
    ppt_app = None
    presentation = None

    try:
        if not os.path.exists(input_ppt_path):
            return False, "输入文件不存在"
        if not input_ppt_path.lower().endswith(('.ppt', '.pptx')):
            return False, "请选择 .ppt 或 .pptx 文件"

        status_callback("正在初始化PowerPoint...")
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = True   # 必须可见，避免权限错误

        presentation = ppt_app.Presentations.Open(input_ppt_path, WithWindow=False)
        status_callback("已打开演示文稿，正在分析页面尺寸...")

        slide_width_points = presentation.PageSetup.SlideWidth
        slide_height_points = presentation.PageSetup.SlideHeight

        # 计算目标像素尺寸
        width_px = int(round(slide_width_points / 72.0 * dpi_target))
        height_px = int(round(slide_height_points / 72.0 * dpi_target))

        # 限制过大尺寸（避免内存溢出）
        max_px = 16000
        if width_px > max_px:
            width_px = max_px
            height_px = int(round(width_px * slide_height_points / slide_width_points))
        if height_px > max_px:
            height_px = max_px
            width_px = int(round(height_px * slide_width_points / slide_height_points))

        status_callback(f"目标图片尺寸：{width_px} x {height_px} 像素 (DPI={dpi_target})")

        temp_dir = tempfile.mkdtemp(prefix="ppt_highres_")
        slides = []
        for slide in presentation.Slides:
            if slide.SlideShowTransition.Hidden:
                continue
            slides.append(slide)

        total_slides = len(slides)
        if total_slides == 0:
            return False, "未找到任何非隐藏幻灯片"

        # 第一步：导出高分辨率图片（PNG）
        png_files = []
        for idx, slide in enumerate(slides, start=1):
            progress_callback(idx, total_slides, f"正在导出第 {idx}/{total_slides} 页...")
            png_filename = f"slide_{idx:04d}.png"
            png_path = os.path.join(temp_dir, png_filename)
            slide.Export(png_path, "PNG", width_px, height_px)
            # 等待文件写入完成
            for _ in range(30):
                if os.path.exists(png_path) and os.path.getsize(png_path) > 0:
                    break
                import time
                time.sleep(0.1)
            png_files.append(png_path)

        status_callback("图片导出完成，正在处理图片...")

        # 关闭并清理PowerPoint
        presentation.Close()
        ppt_app.Quit()
        ppt_app = None

        # 第二步：如果需要压缩，将PNG转为JPEG（质量可调）
        image_files = []  # 存储最终用于插入的图片路径
        if enable_compress:
            status_callback(f"正在压缩图片（质量 {compress_quality}%）...")
            for idx, png_path in enumerate(png_files, start=1):
                progress_callback(idx, len(png_files), f"压缩第 {idx}/{len(png_files)} 页图片...")
                jpg_path = png_path.replace('.png', '.jpg')
                with Image.open(png_path) as img:
                    # 如果图片模式为RGBA，转换为RGB（JPEG不支持透明通道）
                    if img.mode == 'RGBA':
                        rgb_img = Image.new('RGB', img.size, (255, 255, 255))
                        rgb_img.paste(img, mask=img.split()[3])  # 使用alpha通道作为mask
                        rgb_img.save(jpg_path, 'JPEG', quality=compress_quality, optimize=True)
                    else:
                        img.save(jpg_path, 'JPEG', quality=compress_quality, optimize=True)
                # 删除原始PNG文件，释放空间
                os.remove(png_path)
                image_files.append(jpg_path)
        else:
            # 不压缩，直接使用PNG
            image_files = png_files

        status_callback("正在生成图片型PPT...")

        # 第三步：创建新的PPTX，插入图片
        new_prs = PPTXPresentation()
        new_prs.slide_width = Emu(int(slide_width_points * 12700))
        new_prs.slide_height = Emu(int(slide_height_points * 12700))

        try:
            blank_layout = new_prs.slide_layouts[6]
        except IndexError:
            blank_layout = new_prs.slide_layouts[0]

        # 确保图片按页码顺序排列
        image_files.sort(key=lambda x: int(os.path.basename(x).split('_')[1].split('.')[0]))

        for i, img_path in enumerate(image_files, start=1):
            progress_callback(i, len(image_files), f"正在插入第 {i}/{len(image_files)} 页图片...")
            slide = new_prs.slides.add_slide(blank_layout)

            # 清除默认占位符
            for shape in list(slide.shapes):
                sp = shape.element
                sp.getparent().remove(sp)

            # 设置白色背景
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)

            # 获取图片实际像素尺寸
            with Image.open(img_path) as pil_img:
                img_px_w, img_px_h = pil_img.size

            # 计算图片物理尺寸（点）
            img_pt_w = img_px_w * 72.0 / dpi_target
            img_pt_h = img_px_h * 72.0 / dpi_target

            # 缩放以完整显示在幻灯片内
            scale = min(slide_width_points / img_pt_w, slide_height_points / img_pt_h)
            display_pt_w = img_pt_w * scale
            display_pt_h = img_pt_h * scale

            left_pt = (slide_width_points - display_pt_w) / 2.0
            top_pt = (slide_height_points - display_pt_h) / 2.0

            left_emu = Emu(int(left_pt * 12700))
            top_emu = Emu(int(top_pt * 12700))
            width_emu = Emu(int(display_pt_w * 12700))
            height_emu = Emu(int(display_pt_h * 12700))

            slide.shapes.add_picture(img_path, left_emu, top_emu, width_emu, height_emu)

        new_prs.save(output_pptx_path)
        status_callback("图片型PPT已生成完毕！")
        return True, None

    except Exception as e:
        error_msg = f"转换失败: {str(e)}"
        return False, error_msg

    finally:
        if presentation:
            try:
                presentation.Close()
            except:
                pass
        if ppt_app:
            try:
                ppt_app.Quit()
            except:
                pass
        if temp_dir and os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
            except:
                pass


# ========== GUI界面（新增DPI选择和压缩选项） ==========
class App:
    def __init__(self, root):
        self.root = root
        root.title("PPT 转 高分辨率图片型PPT（支持压缩）")
        root.geometry("650x520")
        root.resizable(False, False)

        self.input_path = tk.StringVar()
        self.output_path = tk.StringVar()
        self.is_running = False

        # DPI相关
        self.dpi_var = tk.StringVar(value="1200")
        self.custom_dpi_var = tk.StringVar()
        self.custom_dpi_entry = None
        self.dpi_combobox = None

        # 压缩相关
        self.compress_var = tk.BooleanVar(value=False)
        self.quality_var = tk.IntVar(value=85)
        self.quality_label = None
        self.quality_scale = None

        self.create_widgets()

    def create_widgets(self):
        row = 0
        # 输入文件
        tk.Label(self.root, text="输入PPT/PPTX文件：").grid(row=row, column=0, padx=10, pady=10, sticky="w")
        entry_input = tk.Entry(self.root, textvariable=self.input_path, width=50)
        entry_input.grid(row=row, column=1, padx=5, pady=10)
        btn_browse_input = tk.Button(self.root, text="浏览...", command=self.browse_input)
        btn_browse_input.grid(row=row, column=2, padx=5, pady=10)
        row += 1

        # 输出文件
        tk.Label(self.root, text="输出PPTX文件：").grid(row=row, column=0, padx=10, pady=10, sticky="w")
        entry_output = tk.Entry(self.root, textvariable=self.output_path, width=50)
        entry_output.grid(row=row, column=1, padx=5, pady=10)
        btn_browse_output = tk.Button(self.root, text="保存为...", command=self.browse_output)
        btn_browse_output.grid(row=row, column=2, padx=5, pady=10)
        row += 1

        # DPI选择区域
        dpi_frame = tk.LabelFrame(self.root, text="导出图片DPI (≥72)", padx=5, pady=5)
        dpi_frame.grid(row=row, column=0, columnspan=3, padx=10, pady=5, sticky="ew")
        tk.Label(dpi_frame, text="常用DPI：").pack(side=tk.LEFT, padx=5)
        self.dpi_combobox = ttk.Combobox(dpi_frame, textvariable=self.dpi_var, values=["96", "150", "300", "600", "1200", "自定义"], width=10)
        self.dpi_combobox.pack(side=tk.LEFT, padx=5)
        self.dpi_combobox.bind("<<ComboboxSelected>>", self.on_dpi_selected)
        tk.Label(dpi_frame, text="自定义DPI：").pack(side=tk.LEFT, padx=(20,5))
        self.custom_dpi_entry = tk.Entry(dpi_frame, textvariable=self.custom_dpi_var, width=8, state="disabled")
        self.custom_dpi_entry.pack(side=tk.LEFT, padx=5)
        row += 1

        # 压缩选项区域
        compress_frame = tk.LabelFrame(self.root, text="图片压缩（仅当选择压缩时有效）", padx=5, pady=5)
        compress_frame.grid(row=row, column=0, columnspan=3, padx=10, pady=5, sticky="ew")
        cb_compress = tk.Checkbutton(compress_frame, text="启用压缩（大幅减小PPT文件体积）", variable=self.compress_var, command=self.on_compress_toggle)
        cb_compress.pack(anchor="w", padx=5)
        # 质量滑块
        inner_frame = tk.Frame(compress_frame)
        inner_frame.pack(fill="x", padx=10, pady=5)
        tk.Label(inner_frame, text="JPEG质量：").pack(side=tk.LEFT)
        self.quality_scale = tk.Scale(inner_frame, from_=1, to=100, orient=tk.HORIZONTAL, variable=self.quality_var, length=200)
        self.quality_scale.pack(side=tk.LEFT, padx=5)
        self.quality_label = tk.Label(inner_frame, text="85%")
        self.quality_label.pack(side=tk.LEFT, padx=5)
        # 质量值实时更新
        self.quality_scale.config(command=self.update_quality_label)
        # 初始状态：压缩未启用时禁用滑块
        self.quality_scale.config(state="disabled")
        self.quality_label.config(state="disabled")
        row += 1

        # 状态栏
        self.status_var = tk.StringVar(value="就绪")
        status_label = tk.Label(self.root, textvariable=self.status_var, relief="sunken", anchor="w")
        status_label.grid(row=row, column=0, columnspan=3, padx=10, pady=5, sticky="ew")
        row += 1

        # 进度条
        self.progress = ttk.Progressbar(self.root, orient="horizontal", length=500, mode="determinate")
        self.progress.grid(row=row, column=0, columnspan=3, padx=10, pady=10)
        row += 1

        self.progress_text = tk.StringVar(value="")
        progress_label = tk.Label(self.root, textvariable=self.progress_text)
        progress_label.grid(row=row, column=0, columnspan=3)
        row += 1

        # 转换按钮
        self.btn_convert = tk.Button(self.root, text="开始转换", command=self.start_conversion, bg="#4CAF50", fg="white", font=("Arial", 10, "bold"))
        self.btn_convert.grid(row=row, column=0, columnspan=3, pady=20)
        row += 1

        # 说明文字
        info_text = """说明：
1. 将每一页PPT导出为指定DPI的高清图片，重新组合成图片型PPT（不可编辑文字）。
2. 启用压缩会将图片转为JPEG格式，可大幅减小文件体积，但会略微损失画质（质量85%通常视觉无损）。
3. 仅支持Windows系统，需要安装Microsoft PowerPoint。
4. 转换过程中PowerPoint窗口会短暂弹出，属正常现象。"""
        info_label = tk.Label(self.root, text=info_text, justify="left", fg="gray", font=("Arial", 9))
        info_label.grid(row=row, column=0, columnspan=3, padx=10, pady=10, sticky="w")

    def on_dpi_selected(self, event=None):
        if self.dpi_var.get() == "自定义":
            self.custom_dpi_entry.config(state="normal")
            self.custom_dpi_entry.focus()
        else:
            self.custom_dpi_entry.config(state="disabled")
            self.custom_dpi_var.set("")

    def on_compress_toggle(self):
        if self.compress_var.get():
            self.quality_scale.config(state="normal")
            self.quality_label.config(state="normal")
        else:
            self.quality_scale.config(state="disabled")
            self.quality_label.config(state="disabled")

    def update_quality_label(self, val):
        self.quality_label.config(text=f"{int(float(val))}%")

    def get_effective_dpi(self):
        """获取用户选择的DPI（整数）"""
        dpi_str = self.dpi_var.get()
        if dpi_str == "自定义":
            try:
                custom = int(self.custom_dpi_var.get())
                if custom < 72:
                    raise ValueError
                return custom
            except:
                messagebox.showerror("错误", "自定义DPI必须是≥72的整数")
                return None
        else:
            return int(dpi_str)

    def browse_input(self):
        file_path = filedialog.askopenfilename(
            title="选择PPT/PPTX文件",
            filetypes=[("PowerPoint文件", "*.ppt *.pptx"), ("所有文件", "*.*")]
        )
        if file_path:
            self.input_path.set(file_path)
            # 自动建议输出路径
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            out_dir = os.path.dirname(file_path)
            default_out = os.path.join(out_dir, f"{base_name}_图片型.pptx")
            self.output_path.set(default_out)

    def browse_output(self):
        file_path = filedialog.asksaveasfilename(
            title="保存PPTX文件",
            defaultextension=".pptx",
            filetypes=[("PowerPoint文件", "*.pptx"), ("所有文件", "*.*")]
        )
        if file_path:
            self.output_path.set(file_path)

    def start_conversion(self):
        if self.is_running:
            messagebox.showwarning("提示", "转换任务正在进行中，请稍后再试。")
            return

        input_file = self.input_path.get().strip()
        output_file = self.output_path.get().strip()

        if not input_file:
            messagebox.showerror("错误", "请选择输入PPT/PPTX文件。")
            return
        if not output_file:
            messagebox.showerror("错误", "请指定输出PPTX文件的保存路径。")
            return
        if not output_file.lower().endswith('.pptx'):
            output_file += '.pptx'
            self.output_path.set(output_file)

        dpi = self.get_effective_dpi()
        if dpi is None:
            return

        enable_compress = self.compress_var.get()
        compress_quality = self.quality_var.get() if enable_compress else 85

        # 禁用转换按钮，启动线程
        self.btn_convert.config(state="disabled", text="转换中...")
        self.is_running = True
        self.progress["value"] = 0
        self.progress_text.set("")
        self.status_var.set("正在转换，请稍候...")

        thread = threading.Thread(target=self.conversion_task,
                                  args=(input_file, output_file, dpi, enable_compress, compress_quality),
                                  daemon=True)
        thread.start()

    def conversion_task(self, input_file, output_file, dpi, enable_compress, quality):
        def progress_callback(current, total, msg):
            self.root.after(0, self.update_progress, current, total, msg)

        def status_callback(msg):
            self.root.after(0, self.update_status, msg)

        success, error_msg = ppt_to_high_res_ppt(
            input_file, output_file, dpi, enable_compress, quality,
            progress_callback, status_callback
        )
        self.root.after(0, self.conversion_finished, success, error_msg, output_file)

    def update_progress(self, current, total, msg):
        if total > 0:
            percent = int(current / total * 100)
            self.progress["value"] = percent
        self.progress_text.set(msg)
        self.status_var.set(msg)

    def update_status(self, msg):
        self.status_var.set(msg)

    def conversion_finished(self, success, error_msg, output_file):
        self.is_running = False
        self.btn_convert.config(state="normal", text="开始转换")
        if success:
            messagebox.showinfo("转换成功", f"图片型PPT已生成：\n{output_file}")
            self.status_var.set("转换完成")
            self.progress["value"] = 100
        else:
            messagebox.showerror("转换失败", error_msg)
            self.status_var.set("转换失败")
            self.progress["value"] = 0
            self.progress_text.set("")


if __name__ == "__main__":
    root = tk.Tk()
    app = App(root)
    root.mainloop()